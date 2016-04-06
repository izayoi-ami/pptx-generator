from pptx import Presentation
import argparse
import os 

BLANK_LAYOUT = 0
TITLE_CONTENT_LAYOUT = 1
parser = argparse.ArgumentParser("Powerpoint Generator")
parser.add_argument("outfile")
parser.add_argument("images", nargs="+")
parser.add_argument("-f", "--full", action="store_true", help="Ignore aspect ratio and fill the picture to whole content container")
args = parser.parse_args()
files = args.images
prs = Presentation()
tsl = prs.slide_layouts[TITLE_CONTENT_LAYOUT]
first_slide = prs.slides.add_slide(tsl)
content = first_slide.placeholders[1]
title = first_slide.shapes.title
title.text = args.outfile
dim = ["left","top","width","height"]
get_dim = lambda x:{
        "left" : x.left,
        "top" : x.top,
        "width": x.width,
        "height": x.height
        }
title_dim = get_dim(title)
content_dim = get_dim(content)
if not args.full:
    content_dim.pop("height", None)
        
for f in files:
    slide_layout = prs.slide_layouts[TITLE_CONTENT_LAYOUT]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = os.path.basename(f)
    #title = slide.shapes.add_textbox(**title_dim)
    #title.text = f
    slide.shapes.add_picture(f,**content_dim)
prs.save(args.outfile)
